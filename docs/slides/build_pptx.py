"""Build the 4-slide SDMX Surfer deck as native, unthemed PowerPoint shapes.

Everything is editable: the architecture diagram and the snapshot wireframe are
rebuilt from shapes (no baked images), so a later theming pass can restyle all
of it. 16:9, default fonts, grayscale except the diagram's own light hues.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
INK = RGBColor(0x11, 0x11, 0x11)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHTGRAY = RGBColor(0x99, 0x99, 0x99)
HAIR = RGBColor(0xDD, 0xDD, 0xDD)

BLUE_F, BLUE_S = RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0x1D, 0x4E, 0xD8)
VIOLET_F, VIOLET_S = RGBColor(0xED, 0xE9, 0xFE), RGBColor(0x6D, 0x28, 0xD9)
AMBER_F, AMBER_S = RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xD9, 0x77, 0x06)
GREEN_F, GREEN_S = RGBColor(0xDC, 0xFC, 0xE7), RGBColor(0x15, 0x80, 0x3D)
STONE_F, STONE_S = RGBColor(0xF5, 0xF5, 0xF4), RGBColor(0x44, 0x40, 0x3C)
GREEN_TXT = RGBColor(0x16, 0x65, 0x34)
AMBER_TXT = RGBColor(0xB4, 0x53, 0x09)
BLUE_TXT = RGBColor(0x1D, 0x4E, 0xD8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- helpers

def add_text(slide, x, y, w, h, runs, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, line_spacing=1.15, anchor=MSO_ANCHOR.TOP):
    """runs: str, or list of paragraphs; each paragraph is a list of
    (text, bold, color) run tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, bold, color)]]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        for text, b, c in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = b
            r.font.color.rgb = c
    return tb


def add_bullets(slide, x, y, w, h, items, size=13.5, gap=10):
    """items: list of (lead, rest) — lead is bold."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for lead, rest in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.18
        p.space_after = Pt(gap)
        dash = p.add_run(); dash.text = "–  "
        dash.font.size = Pt(size); dash.font.color.rgb = LIGHTGRAY
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = INK
        if rest:
            r2 = p.add_run(); r2.text = " " + rest
            r2.font.size = Pt(size); r2.font.color.rgb = INK
    return tb


def box(slide, x, y, w, h, fill, stroke, stroke_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = 0.12
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if stroke is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = stroke; s.line.width = Pt(stroke_w)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    return s


def box_label(s, title, sub=None, title_size=12, sub_size=9.5,
              title_color=INK, sub_color=GRAY):
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(title_size); r.font.bold = True; r.font.color.rgb = title_color
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(sub_size); r2.font.color.rgb = sub_color


def arrow(slide, x1, y1, x2, y2, color=INK, weight=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(weight)
    # arrowhead at the end
    ln = c.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'arrow', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    c.shadow.inherit = False
    return c


def footer(slide, left, right):
    pass  # slides are embedded in a larger deck; no per-slide footers


def heading(slide, text, kicker=None):
    y = Inches(0.55)
    if kicker:
        add_text(slide, Inches(0.65), Inches(0.38), Inches(11), Inches(0.3),
                 kicker.upper(), size=10.5, color=GRAY, bold=True)
        y = Inches(0.72)
    add_text(slide, Inches(0.65), y, Inches(12), Inches(0.7),
             text, size=27, bold=True)


# ════════════════════════════════════ Slide 1 — What it is
s = prs.slides.add_slide(BLANK)
heading(s, "Dashboards you ask for, numbers you can trust",
        kicker="SDMX Surfer · Country Snapshots")

add_bullets(s, Inches(0.65), Inches(1.75), Inches(6.4), Inches(4.6), [
    ("SDMX Surfer", "is a web app where you build dashboards on official "
     "statistics by describing them in plain language. An AI agent finds the "
     "right data in the SDMx catalogue and assembles the charts."),
    ("Country Snapshots", "is the curated layer on top: ready-made "
     "statistical profiles for every Pacific country and territory, with "
     "comparison and regional views."),
    ("Both run on the same pipeline:", "every number is fetched live from "
     ".Stat at the moment you look at it, and every chart links back to its "
     "official source."),
], size=14, gap=14)

panel = box(s, Inches(7.5), Inches(1.75), Inches(5.2), Inches(4.3), None, HAIR)
add_text(s, Inches(7.85), Inches(2.0), Inches(4.6), Inches(0.3),
         "THREE PROPERTIES, BY CONSTRUCTION", size=10.5, color=GRAY, bold=True)
badges = [
    "Live data – no copies, no extracts",
    "AI-assisted – conversation, not query languages",
    "Human-verifiable – every chart cites its source",
]
for i, t in enumerate(badges):
    b = box(s, Inches(7.85), Inches(2.45 + i * 0.72), Inches(4.55), Inches(0.52),
            RGBColor(0xFF, 0xFF, 0xFF), INK, 1.5)
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t
    r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = INK
add_text(s, Inches(7.85), Inches(4.85), Inches(4.6), Inches(1.0),
         "The AI cannot touch the data itself: it writes the query, and .Stat "
         "serves the values straight to the chart (next slide).",
         size=11.5, color=GRAY)
footer(s, "sdmxsurfer.net", "1 / 4")

# ════════════════════════════════════ Slide 2 — Architecture
s = prs.slides.add_slide(BLANK)
heading(s, "The AI writes the question, never the numbers")

add_bullets(s, Inches(0.65), Inches(1.6), Inches(5.1), Inches(5.2), [
    ("Describe the dashboard in plain language.", "The agent explores the "
     "statistical catalogue step by step (dataflows, dimensions, availability) "
     "and writes a small, human-readable dashboard spec."),
    ("Data never passes through the model.", "Charts fetch values directly "
     "from the .Stat SDMx API at view time. The AI cannot corrupt, invent, or "
     "round a number: it only chooses the query and the chart type."),
    ("Always live, always checkable.", "No copies or extracts; reload and you "
     "see today's data. Every panel cites its exact source query, one click "
     "from the official .Stat page, and the spec stays human-editable."),
], size=12.5, gap=12)

# --- diagram, rebuilt as native shapes -----------------------------------
DX, DY = Inches(6.05), Inches(1.45)   # diagram origin
def dx(v): return DX + Emu(int(Inches(v)))
def dy(v): return DY + Emu(int(Inches(v)))

# phase bands
band1 = box(s, dx(1.30), dy(0.00), Inches(5.55), Inches(1.10), RGBColor(0xEF, 0xF5, 0xFE), None, shape=MSO_SHAPE.RECTANGLE)
band2 = box(s, dx(1.30), dy(1.10), Inches(5.55), Inches(0.95), RGBColor(0xFE, 0xFA, 0xEC), None, shape=MSO_SHAPE.RECTANGLE)
band3 = box(s, dx(1.30), dy(2.05), Inches(5.55), Inches(1.15), RGBColor(0xEF, 0xFB, 0xF3), None, shape=MSO_SHAPE.RECTANGLE)

add_text(s, dx(0.0), dy(0.02), Inches(1.25), Inches(0.25), "AUTHORING", size=9, color=GRAY, bold=True)
add_text(s, dx(0.0), dy(1.12), Inches(1.45), Inches(0.25), "CONFIGURATION", size=9, color=AMBER_TXT, bold=True)
add_text(s, dx(0.0), dy(2.07), Inches(1.25), Inches(0.25), "RUNTIME", size=9, color=GREEN_TXT, bold=True)

# boxes — authoring row
b_user = box(s, dx(0.10), dy(0.30), Inches(1.05), Inches(0.62), RGBColor(0xFF, 0xFF, 0xFF), INK)
box_label(b_user, "User", "dashboard goal")
b_agent = box(s, dx(1.55), dy(0.30), Inches(1.30), Inches(0.62), BLUE_F, BLUE_S)
box_label(b_agent, "LLM Agent", "query + chart")
b_mcp = box(s, dx(3.30), dy(0.30), Inches(1.35), Inches(0.62), VIOLET_F, VIOLET_S)
box_label(b_mcp, "MCP Gateway", "SDMX discovery")
# endpoint spans all rows
b_ep = box(s, dx(5.35), dy(0.30), Inches(1.30), Inches(2.75), GREEN_F, GREEN_S)
tf = b_ep.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "SDMX\nEndpoint"; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = INK
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "\nmetadata"; r2.font.size = Pt(9.5); r2.font.color.rgb = GREEN_TXT
p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "\n\n\n\n\ndata values"; r3.font.size = Pt(9.5); r3.font.color.rgb = GREEN_TXT

# boxes — configuration row
b_cfg = box(s, dx(1.55), dy(1.25), Inches(1.30), Inches(0.62), AMBER_F, AMBER_S)
box_label(b_cfg, "JSON config", "human-readable")
b_lib = box(s, dx(3.30), dy(1.25), Inches(1.35), Inches(0.62), STONE_F, STONE_S)
box_label(b_lib, "Dashboard", "components")

# boxes — runtime row
b_dash = box(s, dx(3.30), dy(2.30), Inches(1.35), Inches(0.72), RGBColor(0xE7, 0xE5, 0xE4), STONE_S)
box_label(b_dash, "Interactive\ndashboard", "fetches data")

# arrows
arrow(s, dx(1.15), dy(0.61), dx(1.55), dy(0.61), INK)                      # user → agent
arrow(s, dx(2.85), dy(0.50), dx(3.30), dy(0.50), BLUE_S)                   # agent → mcp
arrow(s, dx(3.30), dy(0.74), dx(2.85), dy(0.74), BLUE_S)                   # mcp → agent
arrow(s, dx(4.65), dy(0.50), dx(5.35), dy(0.65), BLUE_S)                   # mcp → endpoint
arrow(s, dx(5.35), dy(0.85), dx(4.65), dy(0.74), BLUE_S)                   # endpoint → mcp
arrow(s, dx(2.20), dy(0.92), dx(2.20), dy(1.25), AMBER_S)                  # agent writes config
arrow(s, dx(2.85), dy(1.56), dx(3.30), dy(1.56), AMBER_S)                  # config → components
arrow(s, dx(3.97), dy(1.87), dx(3.97), dy(2.30), INK)                      # renders
arrow(s, dx(4.65), dy(2.50), dx(5.35), dy(2.50), GREEN_S, 2.0)             # sdmx query
arrow(s, dx(5.35), dy(2.85), dx(4.65), dy(2.85), GREEN_S, 2.0)             # values

# small arrow labels
add_text(s, dx(2.30), dy(0.98), Inches(0.75), Inches(0.22), "writes", size=8.5, color=AMBER_TXT)
add_text(s, dx(4.05), dy(1.95), Inches(0.75), Inches(0.22), "renders", size=8.5, color=GRAY)
add_text(s, dx(4.62), dy(2.24), Inches(0.85), Inches(0.2), "SDMX query", size=8, color=GREEN_TXT)
add_text(s, dx(4.72), dy(2.90), Inches(0.7), Inches(0.2), "values", size=8, color=GREEN_TXT)
add_text(s, dx(2.86), dy(0.24), Inches(0.75), Inches(0.2), "MCP call", size=8, color=BLUE_TXT)

add_text(s, DX, dy(3.35), Inches(6.7), Inches(0.55),
         "The model acts only in the authoring phase. At runtime the "
         "dashboard talks to the SDMx endpoint directly.",
         size=10, color=GRAY)
footer(s, "SDMX Surfer", "2 / 4")

# ════════════════════════════════════ Slide 3 — Country Snapshots
s = prs.slides.add_slide(BLANK)
heading(s, "Country Snapshots: the curated layer")

add_bullets(s, Inches(0.65), Inches(1.6), Inches(5.6), Inches(5.2), [
    ("Built for MFAT, useful to everyone.", "Developed for the New Zealand "
     "Ministry of Foreign Affairs and Trade: one-page statistical profiles "
     "ready to brief from, with a curated catalogue of ~100 indicators "
     "across 12 themes for each of 22 Pacific countries and territories, "
     "plus side-by-side comparison and region-wide views."),
    ("Same trust pipeline.", "Every indicator on the page is a live .Stat "
     "query rendered on load, with its source cited. PDF export keeps the "
     "source links clickable."),
    ("From reading to exploring in one click.", "“Explore in Surfer” turns "
     "any snapshot page into an editable AI session: change countries, time "
     "ranges, or chart types by asking."),
], size=13, gap=14)

# --- wireframe, native shapes --------------------------------------------
WX, WY = Inches(6.7), Inches(1.55)
frame = box(s, WX, WY, Inches(6.0), Inches(4.6), RGBColor(0xFD, 0xFD, 0xFC), HAIR)
hdr = box(s, WX, WY, Inches(6.0), Inches(0.5), RGBColor(0xF2, 0xF2, 0xEF), None, shape=MSO_SHAPE.RECTANGLE)
box(s, WX + Inches(0.25), WY + Inches(0.16), Inches(1.5), Inches(0.17), RGBColor(0xC9, 0xC9, 0xC4), None, shape=MSO_SHAPE.RECTANGLE)
box(s, WX + Inches(4.3), WY + Inches(0.14), Inches(0.72), Inches(0.2), RGBColor(0xDD, 0xDD, 0xD8), None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
box(s, WX + Inches(5.1), WY + Inches(0.14), Inches(0.72), Inches(0.2), RGBColor(0xDD, 0xDD, 0xD8), None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
# title bars
box(s, WX + Inches(0.25), WY + Inches(0.72), Inches(2.3), Inches(0.22), RGBColor(0xB9, 0xB9, 0xB3), None, shape=MSO_SHAPE.RECTANGLE)
box(s, WX + Inches(0.25), WY + Inches(1.02), Inches(1.2), Inches(0.13), RGBColor(0xD9, 0xD9, 0xD4), None, shape=MSO_SHAPE.RECTANGLE)
# TOC rail
box(s, WX + Inches(0.25), WY + Inches(1.3), Inches(1.05), Inches(2.5), RGBColor(0xF5, 0xF5, 0xF2), None)
for i, wds in enumerate([0.8, 0.68, 0.74, 0.6]):
    box(s, WX + Inches(0.37), WY + Inches(1.48 + i * 0.2), Inches(wds), Inches(0.1),
        RGBColor(0xCF, 0xCF, 0xCA) if i == 0 else RGBColor(0xDC, 0xDC, 0xD7),
        None, shape=MSO_SHAPE.RECTANGLE)
# chart cards
card1 = box(s, WX + Inches(1.5), WY + Inches(1.3), Inches(2.1), Inches(1.55), RGBColor(0xFF, 0xFF, 0xFF), HAIR)
box(s, WX + Inches(1.65), WY + Inches(1.45), Inches(1.15), Inches(0.12), RGBColor(0xCF, 0xCF, 0xCA), None, shape=MSO_SHAPE.RECTANGLE)
ln1 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, WX + Inches(1.7), WY + Inches(2.6), WX + Inches(3.4), WY + Inches(1.85))
ln1.line.color.rgb = RGBColor(0x9A, 0x9A, 0x94); ln1.line.width = Pt(2); ln1.shadow.inherit = False
ln2 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, WX + Inches(1.7), WY + Inches(2.68), WX + Inches(3.4), WY + Inches(2.25))
ln2.line.color.rgb = RGBColor(0xC4, 0xC4, 0xBE); ln2.line.width = Pt(2); ln2.shadow.inherit = False
card2 = box(s, WX + Inches(3.75), WY + Inches(1.3), Inches(2.1), Inches(1.55), RGBColor(0xFF, 0xFF, 0xFF), HAIR)
box(s, WX + Inches(3.9), WY + Inches(1.45), Inches(1.15), Inches(0.12), RGBColor(0xCF, 0xCF, 0xCA), None, shape=MSO_SHAPE.RECTANGLE)
bars = [(0.0, 0.55), (0.32, 0.38), (0.64, 0.68), (0.96, 0.27), (1.28, 0.47)]
for bx, bh in bars:
    box(s, WX + Inches(3.95 + bx), WY + Inches(2.72 - bh), Inches(0.24), Inches(bh),
        RGBColor(0xB9, 0xB9, 0xB3), None, shape=MSO_SHAPE.RECTANGLE)
# KPI cards
for cx in (1.5, 3.75):
    box(s, WX + Inches(cx), WY + Inches(3.0), Inches(2.1), Inches(0.8), RGBColor(0xFF, 0xFF, 0xFF), HAIR)
    box(s, WX + Inches(cx + 0.15), WY + Inches(3.13), Inches(1.0), Inches(0.11), RGBColor(0xCF, 0xCF, 0xCA), None, shape=MSO_SHAPE.RECTANGLE)
    box(s, WX + Inches(cx + 0.15), WY + Inches(3.34), Inches(0.7), Inches(0.3), RGBColor(0xA8, 0xA8, 0xA2), None, shape=MSO_SHAPE.RECTANGLE)

add_text(s, WX, WY + Inches(4.7), Inches(6.0), Inches(0.4),
         "Snapshot page layout (wireframe; live screenshots to follow with the theme).",
         size=10, color=GRAY)
footer(s, "Country Snapshots", "3 / 4")

# ════════════════════════════════════ Slide 4 — Reach + status
s = prs.slides.add_slide(BLANK)
heading(s, "Built for the Pacific, connected to the world")

groups = [
    ("HOME",
     [("SPC's Pacific Data Hub (.Stat)", ", first-class support for all 22 "
       "Pacific countries and territories, treated with regional neutrality.")]),
    ("NATIONAL STATISTICS OFFICES",
     [("Stats NZ · Australian Bureau of Statistics", " · Fiji Bureau of "
       "Statistics · Samoa Bureau of Statistics")]),
    ("INTERNATIONAL ORGANISATIONS",
     [("", "OECD · Eurostat · UNICEF · ILO · IMF · ECB · BIS")]),
]
ty = 1.7
for label, entries in groups:
    add_text(s, Inches(0.65), Inches(ty), Inches(5.9), Inches(0.25),
             label, size=10, color=GRAY, bold=True)
    lead, rest = entries[0]
    add_text(s, Inches(0.65), Inches(ty + 0.28), Inches(5.9), Inches(0.8),
             [[(lead, True, INK), (rest, False, INK)]], size=13)
    ty += 1.25
add_text(s, Inches(0.65), Inches(5.55), Inches(5.9), Inches(1.1),
         "One conversation can mix sources: Pacific data alongside the same "
         "indicator from any connected provider, every series still fetched "
         "live from its own endpoint.", size=11.5, color=GRAY)

add_bullets(s, Inches(7.1), Inches(1.7), Inches(5.6), Inches(3.4), [
    ("Live today", "at sdmxsurfer.net (capped pilot), moving to "
     "surfer.pacificdata.org."),
    ("Safe to open up:", "per-user daily limits and a hard budget cap on AI "
     "usage; the data itself is public official statistics."),
    ("Owned by the Pacific:", "copyright The Pacific Community (SPC), source "
     "available under a noncommercial licence; built at the Statistics for "
     "Development Division."),
], size=13, gap=13)
sep = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.1), Inches(5.35), Inches(12.7), Inches(5.35))
sep.line.color.rgb = HAIR; sep.line.width = Pt(1); sep.shadow.inherit = False
add_text(s, Inches(7.1), Inches(5.55), Inches(5.6), Inches(1.0),
         "Everything on screen is reproducible: open any chart's source link "
         "and you are on the official .Stat query that produced it.",
         size=11.5, color=INK)
footer(s, "SDMX Surfer · Country Snapshots", "4 / 4")

out = "sdmx-surfer-slides.pptx"
prs.save(out)
print("wrote", out)
